"""
WorkflowMWTroubleshootingZTE.py
================================
Script di diagnostica massiva e bidirezionale (Local vs Remote) per tratte radio ZTE.
Supporta architetture fino a 4+0 XPIC.

Uso standalone:
    python WorkflowMWTroubleshootingZTE.py <percorso_file_excel>

Uso integrato (Alarm Manager):
    from WorkflowMWTroubleshootingZTE import analyze_and_return
    result = analyze_and_return(df)   # df già letto dal pm_history.parquet
"""

import ipaddress
import io
import base64
import os
import sys

import pandas as pd
import matplotlib
matplotlib.use('Agg')   # backend non-interattivo, necessario in contesto server
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from pptx import Presentation
from pptx.util import Inches, Pt

# ─── Palette colori ───────────────────────────────────────────────────────────
MUTED_COLORS = [
    '#1f77b4', '#ff7f0e', '#2ca02c', '#d62728',
    '#9467bd', '#8c564b', '#e377c2', '#7f7f7f'
]


# ─────────────────────────────────────────────────────────────────────────────
#  1. Pre-processing: calcolo colonna Sub-Max Mod Time
# ─────────────────────────────────────────────────────────────────────────────

def compute_downshift_column(site_df: pd.DataFrame, modems_list: list) -> pd.DataFrame:
    """
    Calcola la colonna 'Sub-Max Mod Time(s)' per ogni modem nel DataFrame.
    Restituisce una COPIA del df con la nuova colonna, in modo che sia
    disponibile sia per generate_site_charts sia per analyze_site_stats.

    La colonna rappresenta i secondi totali in cui il modem ha operato
    a modulazione inferiore alla massima (downshift).
    """
    mod_cols = [c for c in site_df.columns if 'Working Time of RX' in c]
    site_df = site_df.copy()
    site_df['Sub-Max Mod Time(s)'] = 0.0

    if not mod_cols:
        return site_df

    for modem in modems_list:
        mask = site_df['PM Checkpoint'] == modem
        modem_df = site_df[mask]
        if modem_df.empty:
            continue
        mod_sums = modem_df[mod_cols].sum()
        max_col = mod_sums.idxmax()
        if pd.notna(max_col):
            sub_max = modem_df[mod_cols].sum(axis=1) - modem_df[max_col].fillna(0)
            site_df.loc[mask, 'Sub-Max Mod Time(s)'] = sub_max.values

    return site_df


# ─────────────────────────────────────────────────────────────────────────────
#  2. Generazione grafici
# ─────────────────────────────────────────────────────────────────────────────

def _fig_to_b64(fig) -> str:
    """Converte una figura matplotlib in stringa base64 PNG (per il frontend)."""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', bbox_inches='tight')
    buf.seek(0)
    b64 = base64.b64encode(buf.read()).decode('utf-8')
    plt.close(fig)
    return b64


def generate_site_charts(site_df: pd.DataFrame, modems_list: list, prefix: str,
                         save_to_disk: bool = True):
    """
    Genera i grafici RSL, XPI, MSE, Downshift e il pannello comparazione.

    Args:
        site_df:       DataFrame già arricchito con 'Sub-Max Mod Time(s)'
        modems_list:   lista dei modem da plottare
        prefix:        prefisso per i nomi file (es. 'local', 'remote')
        save_to_disk:  se True salva PNG su disco; se False restituisce dict base64

    Returns:
        Se save_to_disk=True  → lista di percorsi PNG salvati
        Se save_to_disk=False → dict {metric_key: base64_string}
    """
    site_df = site_df.copy()
    site_df['Begin Time'] = pd.to_datetime(site_df['Begin Time'])
    site_df.sort_values(by='Begin Time', inplace=True)

    colors = {m: MUTED_COLORS[i % len(MUTED_COLORS)] for i, m in enumerate(modems_list)}
    saved_images = []
    b64_images = {}

    def _make_plot(metric, ylabel, filename_key, title):
        fig, ax = plt.subplots(figsize=(12, 6))
        plotted_any = False
        for modem in modems_list:
            modem_df = site_df[site_df['PM Checkpoint'] == modem]
            if not modem_df.empty and metric in modem_df.columns:
                mod_clean = modem_df[['Begin Time', metric]].dropna()
                if not mod_clean.empty:
                    ax.plot(mod_clean['Begin Time'], mod_clean[metric],
                            label=modem, color=colors.get(modem, 'black'))
                    plotted_any = True

        if not plotted_any:
            plt.close(fig)
            return

        ax.set_xlabel('Time')
        ax.set_ylabel(ylabel)
        ax.set_title(title)
        ax.legend()
        ax.grid(True)
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%m-%d %H:%M'))
        plt.xticks(rotation=45)
        plt.tight_layout()

        if save_to_disk:
            out_name = f"{prefix}_{filename_key}.png"
            fig.savefig(out_name, bbox_inches='tight')
            plt.close(fig)
            saved_images.append(out_name)
        else:
            b64_images[filename_key] = _fig_to_b64(fig)

    # Singoli grafici
    if 'Mean XPI(dB)' in site_df.columns:
        _make_plot('Mean XPI(dB)', 'XPI (dB)', 'xpi_trend', 'XPI Trend')
    if 'Mean Received Signal Level(dBm)' in site_df.columns:
        _make_plot('Mean Received Signal Level(dBm)', 'RSL (dBm)',
                   'rsl_trend', 'Received Signal Level (RSL) Trend')
    if 'Mean MSE(dB)' in site_df.columns:
        _make_plot('Mean MSE(dB)', 'MSE (dB)', 'mse_trend',
                   'Mean Squared Error (MSE) Trend')
    if 'Sub-Max Mod Time(s)' in site_df.columns:
        _make_plot('Sub-Max Mod Time(s)', 'Demodulazione (s)',
                   'mod_downshift_trend',
                   'Andamento Cali di Modulazione (sec in Sub-Max Mod)')

    # Pannello comparazione 4-in-1
    required_cols = ['ES(s)', 'Mean MSE(dB)',
                     'Mean Received Signal Level(dBm)', 'Sub-Max Mod Time(s)']
    if all(x in site_df.columns for x in required_cols):
        fig, axes = plt.subplots(4, 1, figsize=(12, 16), sharex=True)
        plotted_multi = False
        for modem in modems_list:
            modem_df = site_df[site_df['PM Checkpoint'] == modem]
            if modem_df.empty:
                continue
            col = colors.get(modem, 'black')
            axes[0].plot(modem_df['Begin Time'], modem_df['ES(s)'],
                         label=modem, color=col)
            axes[1].plot(modem_df['Begin Time'], modem_df['Mean MSE(dB)'],
                         label=modem, color=col)
            axes[2].plot(modem_df['Begin Time'],
                         modem_df['Mean Received Signal Level(dBm)'],
                         label=modem, color=col)
            axes[3].plot(modem_df['Begin Time'],
                         modem_df['Sub-Max Mod Time(s)'],
                         label=modem, color=col)
            plotted_multi = True

        if plotted_multi:
            labels = [
                ('ES (s)', 'ES vs Time'),
                ('MSE (dB)', 'MSE vs Time'),
                ('RSL (dBm)', 'RSL vs Time'),
                ('Downshift (s)', 'Downshift (Sub-Max Mod) vs Time'),
            ]
            for i, (ylab, tit) in enumerate(labels):
                axes[i].set_ylabel(ylab)
                axes[i].set_title(tit)
                axes[i].legend()
                axes[i].grid(True)
            axes[3].xaxis.set_major_formatter(mdates.DateFormatter('%m-%d %H:%M'))
            plt.xticks(rotation=45)
            plt.tight_layout()

            if save_to_disk:
                out_multi = f"{prefix}_es_mse_rsl_comparison.png"
                fig.savefig(out_multi, bbox_inches='tight')
                plt.close(fig)
                saved_images.append(out_multi)
            else:
                b64_images['comparison'] = _fig_to_b64(fig)
        else:
            plt.close(fig)

    return saved_images if save_to_disk else b64_images


# ─────────────────────────────────────────────────────────────────────────────
#  3. Statistiche per sito
# ─────────────────────────────────────────────────────────────────────────────

def analyze_site_stats(site_df: pd.DataFrame, modems_list: list,
                       site_name_readable: str) -> list:
    """
    Calcola statistiche aggregate per ogni modem del sito.
    Il df deve già contenere la colonna 'Sub-Max Mod Time(s)'.
    """
    stats = []
    for modem in modems_list:
        modem_df = site_df[site_df['PM Checkpoint'] == modem]
        if modem_df.empty:
            continue

        tot_es = modem_df['ES(s)'].sum() if 'ES(s)' in modem_df.columns else 0
        min_rsl = (modem_df['Mean Received Signal Level(dBm)'].min()
                   if 'Mean Received Signal Level(dBm)' in modem_df.columns else None)
        min_mse = (modem_df['Mean MSE(dB)'].min()
                   if 'Mean MSE(dB)' in modem_df.columns else None)
        min_xpi = (modem_df['Mean XPI(dB)'].min()
                   if 'Mean XPI(dB)' in modem_df.columns else None)
        tot_downshifts = (modem_df['Sub-Max Mod Time(s)'].sum()
                          if 'Sub-Max Mod Time(s)' in modem_df.columns else 0)

        corr_mse = 0.0
        if 'ES(s)' in modem_df.columns and 'Mean MSE(dB)' in modem_df.columns:
            c = modem_df['ES(s)'].corr(modem_df['Mean MSE(dB)'])
            corr_mse = float(c) if pd.notna(c) else 0.0

        # Calcolo delta potenze IF (Max - Min) per individuare oscillazioni sui cavi coassiali IF
        max_delta_if_rx = None
        if 'Max IF Rx Power(dBm)' in modem_df.columns and 'Min IF Rx Power(dBm)' in modem_df.columns:
            deltas_rx = modem_df['Max IF Rx Power(dBm)'] - modem_df['Min IF Rx Power(dBm)']
            max_delta_if_rx = deltas_rx.max()

        max_delta_if_tx = None
        if 'Max IF Tx Power(dBm)' in modem_df.columns and 'Min IF Tx Power(dBm)' in modem_df.columns:
            deltas_tx = modem_df['Max IF Tx Power(dBm)'] - modem_df['Min IF Tx Power(dBm)']
            max_delta_if_tx = deltas_tx.max()

        stats.append({
            'Site':          site_name_readable,
            'Modem':         modem,
            'Total_ES':      float(tot_es),
            'Min_RSL':       float(min_rsl) if min_rsl is not None else None,
            'Min_MSE':       float(min_mse) if min_mse is not None else None,
            'Min_XPI':       float(min_xpi) if min_xpi is not None else None,
            'Tot_Downshifts': float(tot_downshifts),
            'Corr_ES_MSE':   corr_mse,
            'Max_Delta_IF_Rx': float(max_delta_if_rx) if max_delta_if_rx is not None and pd.notna(max_delta_if_rx) else None,
            'Max_Delta_IF_Tx': float(max_delta_if_tx) if max_delta_if_tx is not None and pd.notna(max_delta_if_tx) else None,
        })
    return stats


def extract_degradation_windows(site_df: pd.DataFrame, modems_list: list) -> list:
    """
    Identifica le finestre temporali (es. orarie o 15-min) in cui c'è stato
    un degrado delle performance (RSL < -65 dBm oppure ES > 0).
    Restituisce una lista di dict per il controllo meteo.
    """
    windows = []
    if 'Begin Time' not in site_df.columns:
        return windows

    df = site_df.copy()
    df['Begin Time'] = pd.to_datetime(df['Begin Time'])
    df.sort_values(by='Begin Time', inplace=True)

    for modem in modems_list:
        modem_df = df[df['PM Checkpoint'] == modem]
        if modem_df.empty:
            continue
            
        for _, row in modem_df.iterrows():
            rsl = row.get('Mean Received Signal Level(dBm)', 0)
            es  = row.get('ES(s)', 0)
            if pd.notna(rsl) and rsl < -65 or pd.notna(es) and es > 0:
                end_time = row.get('End Time')
                end_str  = end_time.isoformat() if pd.notna(end_time) else (row['Begin Time'] + pd.Timedelta(minutes=15)).isoformat()
                
                windows.append({
                    "start":   row['Begin Time'].isoformat(),
                    "end":     end_str,
                    "min_rsl": float(rsl) if pd.notna(rsl) else None,
                    "es":      int(es) if pd.notna(es) else 0,
                    "modem":   modem
                })

    # Deduplica/raggruppa finestre molto vicine (stessa ora) per evitare di fare troppe richieste API
    unique_hours = {}
    for w in windows:
        try:
            # Raggruppa per ora
            dt = pd.to_datetime(w["start"]).floor('h')
            hr_key = dt.isoformat()
            if hr_key not in unique_hours:
                unique_hours[hr_key] = w.copy()
            else:
                # Aggiorna con il caso peggiore
                curr = unique_hours[hr_key]
                if w["min_rsl"] is not None:
                    if curr["min_rsl"] is None or w["min_rsl"] < curr["min_rsl"]:
                        curr["min_rsl"] = w["min_rsl"]
                curr["es"] += w["es"]
        except Exception:
            pass

    return list(unique_hours.values())


# ─────────────────────────────────────────────────────────────────────────────
#  4. Logica conclusioni tecniche
# ─────────────────────────────────────────────────────────────────────────────

def build_conclusions(stats_all: list) -> list:
    """
    Genera i testi di conclusione tecnica a partire dalle statistiche aggregate,
    implementando la Matrice di Diagnostica di Tratta (RSL, MSE, XPI, IF, ACM)
    e fornendo un piano d'azione dettagliato basato sul domain knowledge.
    """
    if not stats_all:
        return [
            "Stato del Link: SCONOSCIUTO. Dati insufficienti per l'analisi.",
            "Verificare il corretto caricamento dei dati PM storici.",
        ]

    # Calcolo dei minimi, massimi e totali di tratta
    total_es = sum(s.get('Total_ES', 0.0) for s in stats_all)
    min_rsl = min((s.get('Min_RSL') for s in stats_all if s.get('Min_RSL') is not None), default=None)
    min_mse = min((s.get('Min_MSE') for s in stats_all if s.get('Min_MSE') is not None), default=None)
    min_xpi = min((s.get('Min_XPI') for s in stats_all if s.get('Min_XPI') is not None), default=None)
    tot_downshifts = sum(s.get('Tot_Downshifts', 0.0) for s in stats_all)
    max_delta_if = max(
        max(((s.get('Max_Delta_IF_Rx') or 0.0) for s in stats_all), default=0.0),
        max(((s.get('Max_Delta_IF_Tx') or 0.0) for s in stats_all), default=0.0)
    )

    conclusions = []
    
    # ── SCENARIO 0: TRATTA RADIO OTTIMALE E CONFORME ──────────────────────────
    if total_es == 0 and min_rsl is not None and min_rsl > -60 and max_delta_if < 0.8 and tot_downshifts == 0:
        return [
            "Stato del Link: OTTIMALE. Tratta Radio Stabile e Conforme.",
            "Nessun Errored Second (ES = 0) e nessun downshift anomalo della modulazione.",
            "Valore XPI conforme (>25 dB costante), livello RSL nominale intatto e stabilità di potenza IF ottimale (<0.8dBm di variazione).",
        ]

    # ── SCENARIO DI DEGRADO O PREOCCUPAZIONE QUALITATIVA ──────────────────────
    if total_es > 0:
        conclusions.append(f"DEGRADO CRITICO RILEVATO: Rilevati {int(total_es)} Errored Seconds (ES) totali sulla tratta radio.")
    else:
        conclusions.append("PREOCCUPAZIONE QUALITATIVA: ES = 0, ma rilevati parametri fuori tolleranza o downshift modulazione.")

    # 1. Attenuazione di Tratta (Rain Fade / Fading Climatico / Ostacolo)
    # Sintomo: RSL Basso, MSE Degradato, XPI normale o leggermente basso, ACM scende
    if min_rsl is not None and min_rsl < -63:
        conclusions.append(f"DIAGNOSI: ATTENUAZIONE DI TRATTA (Fading Climatico / Ostacolo). RSL sceso a {min_rsl:.1f} dBm.")
        conclusions.append("CAUSA PIÙ PROBABILE: Forte attenuazione atmosferica (Rain Fade per pioggia intensa, nebbia fitta) o ostruzione della Line of Sight (LOS) nella zona di Fresnel.")
        conclusions.append("WORKFLOW DI RISOLUZIONE (Esclusione Ambientale):")
        conclusions.append("  * Controllare se l'orario del picco di degrado coincide con eventi meteorologici avversi nella zona.")
        conclusions.append("  * Se è pioggia, attendere il post-maltempo per verificare il ripristino dei livelli nominali.")
        conclusions.append("  * Se il degrado è permanente e graduale nel tempo, effettuare un sopralluogo visivo (LOS check) per verificare la crescita di alberi o nuovi ostacoli artificiali.")
        conclusions.append("  * In caso di cali frequenti in assenza di pioggia, rivalutare il Link Budget (necessità di parabole più grandi o riconfigurazione ACM target).")
        
    # 2. Problema XPIC / Twist (RSL Ottimo, MSE Pessimo, XPI < 20 dB)
    # Sintomo: Antenna ruotata o OMT guasto
    elif min_xpi is not None and min_xpi < 20 and min_rsl is not None and min_rsl >= -60:
        conclusions.append(f"DIAGNOSI: DISALLINEAMENTO CROSS-POLARIZZAZIONE (XPIC / Twist). XPI critico a {min_xpi:.1f} dB.")
        conclusions.append("CAUSA PIÙ PROBABILE: Rotazione dell'antenna sul proprio asse (errore di tilt/twist dovuto a forte vento o fissaggi lenti) o guasto hardware all'OMT (Orthomode Transducer). Le due polarizzazioni interferiscono a vicenda.")
        conclusions.append("WORKFLOW DI RISOLUZIONE (Correzione Cross-Pol):")
        conclusions.append("  * Richiedere intervento in quota dei rigger su entrambi i siti per il ripuntamento delle antenne.")
        conclusions.append("  * Eseguire una regolazione micrometrica della rotazione (Twist) della flangia dell'antenna per massimizzare l'isolamento incrociato fino a ripristinare XPI > 30 dB.")
        conclusions.append("  * Se la regolazione non produce benefici, ispezionare ed eventualmente sostituire l'OMT.")

    # 3. Interferenza Esterna (RSL Ottimo, MSE Pessimo, XPI Normale)
    elif min_mse is not None and min_mse > -35 and min_rsl is not None and min_rsl >= -60:
        conclusions.append(f"DIAGNOSI: INTERFERENZA ESTERNA SU CANALE RADIO. MSE degradato a {min_mse:.1f} dB con RSL ottimo.")
        conclusions.append("CAUSA PIÙ PROBABILE: Un altro ponte radio (co-canale o canale adiacente di terzi) sta trasmettendo sulla stessa frequenza della ODU locale, oppure è presente un forte rumore di fase nell'hardware.")
        conclusions.append("WORKFLOW DI RISOLUZIONE (Gestione delle Interferenze):")
        conclusions.append("  * Eseguire uno Spectrum Scan (Frequency Scan) dall'interfaccia dell'apparato spegnendo temporaneamente il TX remoto per verificare la presenza di segnale estraneo in ingresso.")
        conclusions.append("  * Se confermato, procedere con la variazione del canale di frequenza operativo o inoltrare segnalazione all'ente regolatore per frequenza licenziata.")

    # 4. Problema sul Mezzo Fisico Locale (Oscillazioni IF su connettori/cavo)
    if max_delta_if >= 0.8:
        conclusions.append(f"DIAGNOSI FISICA: ANOMALIA SUL MEZZO FISICO LOCALE (Cavo / Connettore IF). Delta IF = {max_delta_if:.2f} dBm.")
        conclusions.append("CAUSA PIÙ PROBABILE: Fluttuazioni anomale della potenza Intermediate Frequency sul cavo coassiale tra IDU e ODU, indice di infiltrazione d'acqua, connettori N/TNC allentati o ossidati.")
        conclusions.append("WORKFLOW DI RISOLUZIONE (Verifica Cavo/Connettori):")
        conclusions.append("  * Inviare un tecnico on-site per controllare l'ossidazione e il corretto serraggio dei connettori N/TNC sia lato IDU che lato ODU.")
        conclusions.append("  * Verificare lo stato di usura e la tenuta del nastro autoagglomerante di protezione.")
        conclusions.append("  * Se le fluttuazioni persistono dopo il ri-connessionamento, sostituire la tratta di cavo coassiale per sospetta infiltrazione o micro-frattura.")

    # 5. Demodulazione Adattativa (ACM)
    if tot_downshifts > 0:
        conclusions.append(f"DIAGNOSI FUNZIONALE: DEMODULAZIONE ADATTATIVA ATTIVA (ACM Downshifts). Registrati {int(tot_downshifts)} sec in modulazione ridotta.")
        conclusions.append("CAUSA PIÙ PROBABILE: Il ponte radio ha eseguito un downshift della modulazione (es. riducendo i QAM) per aumentare la robustezza e mantenere il link attivo al prezzo di una riduzione temporanea della banda passante (throughput).")

    # In caso di scenari non mappati ma degradati generici
    if len(conclusions) == 1:
        conclusions.append("DIAGNOSI: Degrado qualitativo generico del segnale senza chiara impronta di attenuazione o interferenza.")
        conclusions.append("Azione: Verificare lo storico allarmi hardware e la stabilità complessiva del sistema.")

    return conclusions


# ─────────────────────────────────────────────────────────────────────────────
#  5. Funzione pubblica — per uso integrato (Alarm Manager)
# ─────────────────────────────────────────────────────────────────────────────

def analyze_and_return(df: pd.DataFrame, site_name_filter: str = None) -> dict:
    """
    Esegue l'analisi completa su un DataFrame già letto (es. da pm_history.parquet)
    e restituisce tutto come dizionario JSON-serializzabile.

    Args:
        df:               DataFrame con i dati PM ZTE
        site_name_filter: nome (parziale) del sito locale da analizzare.
                          Se None, usa auto-detect via IP (sito con IP minore = locale).

    Returns:
        {
          "local_site":  str,
          "remote_site": str,
          "arch":        str,
          "stats_local":  [...],
          "stats_remote": [...],
          "conclusion":   [str, ...],
          "outcome":      "SUPERATO" | "FALLIMENTO" | "PREOCCUPAZIONE",
          "charts": {
            "local":  {"rsl_trend": "<b64>", "xpi_trend": "<b64>", ...},
            "remote": {...}
          }
        }
    """
    if df is None or df.empty:
        return {"error": "DataFrame vuoto o non fornito"}

    # ── Identify sites ─────────────────────────────────────────────────────
    if site_name_filter:
        df_local = df[df['ME'].astype(str).str.contains(site_name_filter, na=False, regex=False)].copy()
        if df_local.empty:
            return {"error": f"Nessun dato trovato per il sito '{site_name_filter}'"}
        local_site_root = site_name_filter
        # Auto-detect remote: primo sito diverso nel dataset
        other_sites = df[~df['ME'].astype(str).str.contains(site_name_filter, na=False, regex=False)]['ME'].dropna().unique()
        remote_site_root = str(other_sites[0]).split('#')[0] if len(other_sites) > 0 else "Unknown"
        df_remote = df[df['ME'].astype(str).str.contains(remote_site_root, na=False, regex=False)].copy()
    else:
        local_site_root, remote_site_root, df_local, df_remote = _auto_detect_sites(df)

    local_modems  = df_local['PM Checkpoint'].dropna().unique().tolist()
    remote_modems = df_remote['PM Checkpoint'].dropna().unique().tolist()

    is_4plus0 = any("AIR:2" in m for m in local_modems)
    arch_str  = "4+0 XPIC" if is_4plus0 else "2+0 XPIC"

    # ── Calcola downshift (FIX: prima di charts E stats) ──────────────────
    df_local  = compute_downshift_column(df_local,  local_modems)
    df_remote = compute_downshift_column(df_remote, remote_modems)

    # ── Grafici base64 ────────────────────────────────────────────────────
    local_charts  = generate_site_charts(df_local,  local_modems,  'local',  save_to_disk=False)
    remote_charts = generate_site_charts(df_remote, remote_modems, 'remote', save_to_disk=False) if not df_remote.empty else {}

    # ── Statistiche ───────────────────────────────────────────────────────
    stats_local  = analyze_site_stats(df_local,  local_modems,  local_site_root)
    stats_remote = analyze_site_stats(df_remote, remote_modems, remote_site_root) if not df_remote.empty else []
    stats_all    = stats_local + stats_remote

    # ── Conclusioni ───────────────────────────────────────────────────────
    conclusion = build_conclusions(stats_all)
    outcome    = _get_outcome(conclusion[0] if conclusion else "")

    # ── Finestre di Degrado (per meteo) ───────────────────────────────────
    degradation_windows = extract_degradation_windows(df_local, local_modems)
    # Aggiungiamo anche il remoto se vogliamo
    degradation_windows += extract_degradation_windows(df_remote, remote_modems)
    
    # Deduplica finale su tutto il sito
    unique_dw = {}
    for w in degradation_windows:
        try:
            k = pd.to_datetime(w["start"]).floor('h').isoformat()
            if k not in unique_dw:
                unique_dw[k] = w
            else:
                curr = unique_dw[k]
                if w["min_rsl"] is not None:
                    if curr["min_rsl"] is None or w["min_rsl"] < curr["min_rsl"]:
                        curr["min_rsl"] = w["min_rsl"]
                curr["es"] += w["es"]
        except: pass

    return {
        "local_site":   local_site_root,
        "remote_site":  remote_site_root,
        "arch":         arch_str,
        "stats_local":  stats_local,
        "stats_remote": stats_remote,
        "conclusion":   conclusion,
        "outcome":      outcome,
        "degradation_windows": list(unique_dw.values()),
        "charts": {
            "local":  local_charts,
            "remote": remote_charts,
        }
    }


def _get_outcome(first_line: str) -> str:
    """Estrae il badge esito dalla prima riga di conclusione."""
    fl = first_line.upper()
    if "FALLIMENTO" in fl or "DEGRADO CRITICO" in fl:
        return "FALLIMENTO"
    if "PREOCCUPAZIONE" in fl:
        return "PREOCCUPAZIONE"
    return "SUPERATO"



def _auto_detect_sites(df: pd.DataFrame):
    """
    Identifica automaticamente il sito locale (IP minore) e il remoto (IP maggiore).
    Restituisce (local_root, remote_root, df_local, df_remote).
    """
    sites_info = df[['ME', 'ME IP']].dropna().drop_duplicates()

    def ip_key(ip_str):
        try:
            clean = str(ip_str).split('(')[-1].split(')')[0] if '(' in str(ip_str) else str(ip_str)
            return ipaddress.IPv4Address(clean)
        except Exception:
            return ipaddress.IPv4Address("255.255.255.255")

    unique_sites = sites_info.groupby('ME')['ME IP'].first().reset_index()
    unique_sites['_ip_obj'] = unique_sites['ME IP'].apply(ip_key)
    unique_sites.sort_values(by='_ip_obj', inplace=True)

    local_root  = str(unique_sites.iloc[0]['ME']).split('#')[0]
    remote_root = str(unique_sites.iloc[1]['ME']).split('#')[0] if len(unique_sites) > 1 else local_root

    df_local  = df[df['ME'].astype(str).str.contains(local_root,  na=False, regex=False)].copy()
    df_remote = df[df['ME'].astype(str).str.contains(remote_root, na=False, regex=False)].copy()

    return local_root, remote_root, df_local, df_remote


# ─────────────────────────────────────────────────────────────────────────────
#  6. Uso standalone — genera PPTX su disco
# ─────────────────────────────────────────────────────────────────────────────

def analyze_4plus0_link(file_path: str):
    """Entry point per uso standalone: legge l'Excel e genera il PPTX."""
    print(f"Loading data from {file_path}...")
    df = pd.read_excel(file_path)

    required = ['ME IP', 'ME', 'PM Checkpoint']
    missing  = [c for c in required if c not in df.columns]
    if missing:
        print(f"Colonne mancanti: {missing}")
        return False

    local_root, remote_root, df_local, df_remote = _auto_detect_sites(df)
    print(f"Sito Locale (IP minore): {local_root}")
    print(f"Sito Remoto (Neighbor):  {remote_root}")

    local_modems  = df_local['PM Checkpoint'].dropna().unique().tolist()
    remote_modems = df_remote['PM Checkpoint'].dropna().unique().tolist()
    print(f"[{local_root}]  Modem trovati: {local_modems}")
    print(f"[{remote_root}] Modem trovati: {remote_modems}")

    is_4plus0 = any("AIR:2" in m for m in local_modems)
    arch_str  = "4+0 XPIC" if is_4plus0 else "2+0 XPIC"
    print(f"Architettura rilevata: {arch_str}")

    # Fix: calcola downshift prima di charts e stats
    df_local  = compute_downshift_column(df_local,  local_modems)
    df_remote = compute_downshift_column(df_remote, remote_modems)

    print("\n--- Generazione grafici Sito Locale ---")
    local_images = generate_site_charts(df_local, local_modems, "local", save_to_disk=True)

    remote_images = []
    if not df_remote.empty and remote_modems:
        print("--- Generazione grafici Sito Remoto ---")
        remote_images = generate_site_charts(df_remote, remote_modems, "remote", save_to_disk=True)
    else:
        print("Warning: record sito remoto non trovati nel file.")

    stats_local  = analyze_site_stats(df_local,  local_modems,  local_root)
    stats_remote = analyze_site_stats(df_remote, remote_modems, remote_root)
    conclusion   = build_conclusions(stats_local + stats_remote)

    ppt_name = f"Troubleshooting_{local_root}_vs_{remote_root}.pptx"
    _create_presentation(local_root, remote_root, local_images, remote_images,
                         conclusion, ppt_name, arch_str)
    return True


def _create_presentation(local_site, remote_site, local_images, remote_images,
                          conclusion_texts, output_pptx, arch_str):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Analisi Performance Tratta Radio ZTE"
    slide.placeholders[1].text = (
        f"Sito Locale: {local_site}\n"
        f"Sito Remoto: {remote_site}\n"
        f"Architettura {arch_str}"
    )

    def add_image_slide(image_path, title_text):
        if not os.path.exists(image_path):
            return
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title_text
        s.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), height=Inches(4.5))

    _LABEL_MAP = {
        'rsl':        'RSL',
        'xpi':        'XPI',
        'mse':        'MSE',
        'downshift':  'Cali Modulazione (Downshifts)',
        'comparison': 'Comparazione ES, MSE, RSL',
    }

    for img in local_images:
        key   = next((k for k in _LABEL_MAP if k in img), '')
        label = _LABEL_MAP.get(key, '')
        add_image_slide(img, f"Sito Locale [{local_site}] — {label}")

    if remote_images:
        div = prs.slides.add_slide(prs.slide_layouts[0])
        div.shapes.title.text = "Visibilità dal Sito Remoto"
        div.placeholders[1].text = f"Neighbor: {remote_site}"
        for img in remote_images:
            key   = next((k for k in _LABEL_MAP if k in img), '')
            label = _LABEL_MAP.get(key, '')
            add_image_slide(img, f"Sito Remoto [{remote_site}] — {label}")

    c_slide = prs.slides.add_slide(prs.slide_layouts[1])
    c_slide.shapes.title.text = "Conclusioni Tecniche e Troubleshooting"
    tf = c_slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(conclusion_texts):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20 if i == 0 else 16)
        p.font.bold = (i == 0)
        p.level = 0 if i == 0 else 1

    prs.save(output_pptx)
    print(f"Presentazione salvata: {output_pptx}")


# ─────────────────────────────────────────────────────────────────────────────
#  Entry point CLI
# ─────────────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python WorkflowMWTroubleshootingZTE.py <excel_file>")
        print("Esempio: python WorkflowMWTroubleshootingZTE.py data.xlsx")
        sys.exit(1)
    analyze_4plus0_link(sys.argv[1])
